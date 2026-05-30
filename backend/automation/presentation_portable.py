from pathlib import Path
from copy import deepcopy

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from .lyrics_domain import SlideChunk, Song


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
COLOR_BG = RGBColor(9, 17, 29)
COLOR_BG_ALT = RGBColor(16, 25, 42)
COLOR_PANEL = RGBColor(24, 35, 56)
COLOR_TEXT = RGBColor(244, 247, 252)
COLOR_MUTED = RGBColor(178, 192, 210)
COLOR_TEAL = RGBColor(103, 240, 216)
COLOR_ORANGE = RGBColor(255, 160, 101)


def create_lyrics_presentation(
    output_path: Path,
    songs: list[Song],
    include_welcome_slide: bool,
    include_verse_labels: bool,
    build_song_chunks,
    template_ppt_path: Path | None = None,
) -> None:
    if template_ppt_path is not None:
        presentation = Presentation(str(template_ppt_path))
        _create_template_based_lyrics_presentation(
            presentation=presentation,
            songs=songs,
            include_welcome_slide=include_welcome_slide,
            include_verse_labels=include_verse_labels,
            build_song_chunks=build_song_chunks,
        )
        presentation.save(str(output_path))
        return

    presentation = _base_presentation()

    if include_welcome_slide:
        _add_cover_slide(
            presentation,
            eyebrow="SlideGen",
            title="Lyrics Deck",
            subtitle="Built with the portable presentation engine.",
            accent_color=COLOR_TEAL,
        )

    for song in songs:
        _add_title_slide(
            presentation,
            title=song.title,
            subtitle="Lyrics deck",
            accent_color=COLOR_ORANGE,
        )

        for chunk in build_song_chunks(song, include_verse_labels=include_verse_labels):
            _add_lyrics_slide(presentation, song.title, chunk)

    presentation.save(str(output_path))


def create_verses_presentation(
    output_path: Path,
    book: str,
    chapter: int,
    left_label: str,
    right_label: str,
    left_verses: dict[int, str],
    right_verses: dict[int, str],
    verse_numbers: list[int],
    template_ppt_path: Path | None = None,
    background_image_path: Path | None = None,
    header_font_name: str | None = None,
    verse_font_name: str | None = None,
    header_font_color: str | None = None,
    verse_font_color: str | None = None,
) -> None:
    reference_range = _build_reference_range(book, chapter, verse_numbers)

    if template_ppt_path is not None:
        presentation = Presentation(str(template_ppt_path))
        if _looks_like_structured_verse_template(presentation):
            _create_template_based_verses_presentation(
                presentation=presentation,
                reference_range=reference_range,
                book=book,
                chapter=chapter,
                left_label=left_label,
                right_label=right_label,
                left_verses=left_verses,
                right_verses=right_verses,
                verse_numbers=verse_numbers,
                background_image_path=background_image_path,
                header_font_name=header_font_name,
                verse_font_name=verse_font_name,
                header_font_color=header_font_color,
                verse_font_color=verse_font_color,
            )
            presentation.save(str(output_path))
            return

        presentation = _base_presentation(template_ppt_path=template_ppt_path)
    else:
        presentation = _base_presentation(template_ppt_path=None)

    _add_cover_slide(
        presentation,
        eyebrow="SlideGen",
        title=reference_range,
        subtitle=f"{left_label} and {right_label}",
        accent_color=COLOR_TEAL,
        template_ppt_path=template_ppt_path,
        background_image_path=background_image_path,
        header_font_name=header_font_name,
        header_font_color=header_font_color,
    )

    for verse_number in verse_numbers:
        _add_verse_slide(
            presentation,
            reference=f"{book} {chapter}:{verse_number}",
            left_label=left_label,
            right_label=right_label,
            left_text=left_verses[verse_number],
            right_text=right_verses[verse_number],
            template_ppt_path=template_ppt_path,
            background_image_path=background_image_path,
            header_font_name=header_font_name,
            verse_font_name=verse_font_name,
            header_font_color=header_font_color,
            verse_font_color=verse_font_color,
        )

    presentation.save(str(output_path))


def _base_presentation(template_ppt_path: Path | None = None) -> Presentation:
    presentation = Presentation(str(template_ppt_path)) if template_ppt_path is not None else Presentation()
    if template_ppt_path is None:
        presentation.slide_width = SLIDE_WIDTH
        presentation.slide_height = SLIDE_HEIGHT
    _clear_existing_slides(presentation)
    return presentation


def _build_reference_range(book: str, chapter: int, verse_numbers: list[int]) -> str:
    if not verse_numbers:
        return f"{book} {chapter}"

    start_verse = verse_numbers[0]
    end_verse = verse_numbers[-1]
    if start_verse == end_verse:
        return f"{book} {chapter}:{start_verse}"
    return f"{book} {chapter}:{start_verse}-{end_verse}"


def _create_template_based_lyrics_presentation(
    presentation: Presentation,
    songs: list[Song],
    include_welcome_slide: bool,
    include_verse_labels: bool,
    build_song_chunks,
) -> None:
    original_slide_count = len(presentation.slides)
    catalog = _build_lyrics_prototype_catalog(presentation)

    for song in songs:
        title_slide = _duplicate_slide(presentation, catalog["title_index"])
        _render_template_title_slide(title_slide, song.title)

        for chunk in build_song_chunks(song, include_verse_labels=include_verse_labels):
            lyric_slide = _duplicate_slide(presentation, _choose_lyrics_prototype_index(catalog, chunk))
            _render_template_lyric_slide(lyric_slide, chunk)

    _remove_original_template_slides(
        presentation,
        original_slide_count=original_slide_count,
        keep_first_slide=include_welcome_slide,
    )


def _build_lyrics_prototype_catalog(presentation: Presentation) -> dict[str, object]:
    catalog: dict[str, object] = {}
    lyric_candidates: dict[tuple[int, bool], int] = {}

    for index, slide in enumerate(presentation.slides):
        main_lines, has_label, title_like = _lyrics_prototype_signature(slide)
        if title_like:
            title_text = _normalize_template_text(_primary_text_shape(slide).text)
            if title_text != "WELCOME" and "title_index" not in catalog:
                catalog["title_index"] = index
        if main_lines >= 1 and not title_like:
            lyric_candidates.setdefault((main_lines, has_label), index)

    if "title_index" not in catalog:
        raise ValueError("Could not find a lyrics title-slide prototype in the bundled PowerPoint template.")
    if not lyric_candidates:
        raise ValueError("Could not find lyrics slide prototypes in the bundled PowerPoint template.")

    catalog["lyric_candidates"] = lyric_candidates
    return catalog


def _choose_lyrics_prototype_index(catalog: dict[str, object], chunk: SlideChunk) -> int:
    line_count = len(chunk.lines)
    has_label = bool(chunk.section_label)
    preferred_keys = [
        (line_count, has_label),
        (line_count, False),
        (line_count, True),
    ]

    if line_count == 1:
        preferred_keys.extend([(2, True), (2, False)])
    if line_count == 2:
        preferred_keys.extend([(3, False), (3, True)])
    if line_count == 3:
        preferred_keys.extend([(2, has_label), (4, has_label)])
    if line_count >= 4:
        preferred_keys.extend([(4, True), (4, False), (3, True), (3, False)])

    lyric_candidates = catalog["lyric_candidates"]
    assert isinstance(lyric_candidates, dict)

    seen: set[tuple[int, bool]] = set()
    for key in preferred_keys:
        if key in seen:
            continue
        seen.add(key)
        if key in lyric_candidates:
            return lyric_candidates[key]

    return next(iter(lyric_candidates.values()))


def _lyrics_prototype_signature(slide) -> tuple[int, bool, bool]:
    shapes = _text_shapes(slide)
    if not shapes:
        return (0, False, False)

    main_shape = _primary_text_shape(slide)
    main_lines = len(_shape_text_lines(main_shape))
    extra_shape = _secondary_text_shape(slide)
    has_label = False
    if extra_shape is not None:
        has_label = _normalize_template_text(extra_shape.text).endswith(":")

    title_like = len(shapes) == 1 and main_lines == 1
    return (main_lines, has_label, title_like)


def _render_template_title_slide(slide, song_title: str) -> None:
    title_shape = _primary_text_shape(slide)
    _replace_text_preserving_template_format(title_shape, song_title.upper())


def _render_template_lyric_slide(slide, chunk: SlideChunk) -> None:
    lyric_shape = _primary_text_shape(slide)
    section_shape = _secondary_text_shape(slide)
    _replace_text_preserving_template_format(lyric_shape, "\n".join(chunk.lines))

    if section_shape is not None:
        _replace_text_preserving_template_format(section_shape, chunk.section_label or "")


def _text_shapes(slide) -> list:
    return [shape for shape in slide.shapes if hasattr(shape, "text_frame")]


def _primary_text_shape(slide):
    candidates = []
    for shape in _text_shapes(slide):
        area = int(shape.width) * int(shape.height)
        text_length = len(_normalize_template_text(shape.text))
        candidates.append((area, text_length, shape))

    if not candidates:
        raise ValueError("No text shapes found on lyrics template slide.")

    candidates.sort(key=lambda item: (item[0], item[1]), reverse=True)
    return candidates[0][2]


def _secondary_text_shape(slide):
    candidates = []
    for shape in _text_shapes(slide):
        area = int(shape.width) * int(shape.height)
        candidates.append((area, int(shape.top), int(shape.left), -len(_normalize_template_text(shape.text)), shape))

    if len(candidates) < 2:
        return None

    candidates.sort(key=lambda item: (item[0], item[1], item[2], item[3]))
    return candidates[0][4]


def _shape_text_lines(shape) -> list[str]:
    return [
        _normalize_template_text(paragraph.text)
        for paragraph in shape.text_frame.paragraphs
        if _normalize_template_text(paragraph.text)
    ]


def _normalize_template_text(text: str) -> str:
    return " ".join((text or "").replace("\r", " ").replace("\n", " ").split())


def _looks_like_structured_verse_template(presentation: Presentation) -> bool:
    if len(presentation.slides) < 2:
        return False

    cover_shape_names = {shape.name for shape in presentation.slides[0].shapes}
    verse_shape_names = {shape.name for shape in presentation.slides[1].shapes}
    required_cover = {"TextBox 1", "TextBox 2"}
    required_verse = {"TextBox 2", "TextBox 11", "TextBox 12", "TextBox 13", "TextBox 14"}
    return required_cover.issubset(cover_shape_names) and required_verse.issubset(verse_shape_names)


def _create_template_based_verses_presentation(
    presentation: Presentation,
    reference_range: str,
    book: str,
    chapter: int,
    left_label: str,
    right_label: str,
    left_verses: dict[int, str],
    right_verses: dict[int, str],
    verse_numbers: list[int],
    background_image_path: Path | None = None,
    header_font_name: str | None = None,
    verse_font_name: str | None = None,
    header_font_color: str | None = None,
    verse_font_color: str | None = None,
) -> None:
    _trim_to_slide_count(presentation, 2)
    cover_slide = presentation.slides[0]
    verse_template_slide = presentation.slides[1]

    if background_image_path is not None:
        _insert_background_picture(cover_slide, presentation, background_image_path)
        _insert_background_picture(verse_template_slide, presentation, background_image_path)

    _set_shape_text(cover_slide, "TextBox 2", reference_range)
    _set_shape_text(cover_slide, "TextBox 1", f"{left_label}\nx\n{right_label}")
    if header_font_name:
        _apply_font_override_to_shape(cover_slide, "TextBox 2", header_font_name)
        _apply_font_override_to_shape(cover_slide, "TextBox 1", header_font_name)
    if header_font_color:
        _apply_color_override_to_shape(cover_slide, "TextBox 2", header_font_color)
        _apply_color_override_to_shape(cover_slide, "TextBox 1", header_font_color)

    for index, verse_number in enumerate(verse_numbers):
        slide = verse_template_slide if index == 0 else _duplicate_slide(presentation, 1)
        if index > 0 and background_image_path is not None:
            _insert_background_picture(slide, presentation, background_image_path)
        _set_shape_text(slide, "TextBox 2", f"{book} {chapter}:{verse_number}")
        _set_shape_text(slide, "TextBox 11", left_verses[verse_number])
        _set_shape_text(slide, "TextBox 12", left_label)
        _set_shape_text(slide, "TextBox 13", right_verses[verse_number])
        _set_shape_text(slide, "TextBox 14", right_label)
        if header_font_name:
            _apply_font_override_to_shape(slide, "TextBox 2", header_font_name)
            _apply_font_override_to_shape(slide, "TextBox 12", header_font_name)
            _apply_font_override_to_shape(slide, "TextBox 14", header_font_name)
        if header_font_color:
            _apply_color_override_to_shape(slide, "TextBox 2", header_font_color)
            _apply_color_override_to_shape(slide, "TextBox 12", header_font_color)
            _apply_color_override_to_shape(slide, "TextBox 14", header_font_color)
        if verse_font_name:
            _apply_font_override_to_shape(slide, "TextBox 11", verse_font_name)
            _apply_font_override_to_shape(slide, "TextBox 13", verse_font_name)
        if verse_font_color:
            _apply_color_override_to_shape(slide, "TextBox 11", verse_font_color)
            _apply_color_override_to_shape(slide, "TextBox 13", verse_font_color)


def _trim_to_slide_count(presentation: Presentation, keep_count: int) -> None:
    while len(presentation.slides) > keep_count:
        _delete_slide_at_index(presentation, len(presentation.slides) - 1)


def _remove_original_template_slides(
    presentation: Presentation,
    original_slide_count: int,
    keep_first_slide: bool,
) -> None:
    first_delete_index = 1 if keep_first_slide else 0
    for index in range(original_slide_count - 1, first_delete_index - 1, -1):
        _delete_slide_at_index(presentation, index)


def _delete_slide_at_index(presentation: Presentation, index: int) -> None:
    slide_id = presentation.slides._sldIdLst[index]
    relationship_id = slide_id.rId
    presentation.part.drop_rel(relationship_id)
    del presentation.slides._sldIdLst[index]


def _duplicate_slide(presentation: Presentation, source_index: int):
    source_slide = presentation.slides[source_index]
    duplicated_slide = presentation.slides.add_slide(source_slide.slide_layout)
    relationship_id_map: dict[str, str] = {}

    for shape in list(duplicated_slide.shapes):
        shape.element.getparent().remove(shape.element)

    for shape in source_slide.shapes:
        new_element = deepcopy(shape.element)
        duplicated_slide.shapes._spTree.insert_element_before(new_element, "p:extLst")

    for relation in source_slide.part.rels.values():
        if relation.reltype in {RT.SLIDE_LAYOUT, RT.NOTES_SLIDE}:
            continue
        if relation.is_external:
            relationship_id_map[relation.rId] = duplicated_slide.part.rels.get_or_add_ext_rel(relation.reltype, relation.target_ref)
        else:
            relationship_id_map[relation.rId] = duplicated_slide.part.rels.get_or_add(relation.reltype, relation._target)

    _copy_slide_background(source_slide, duplicated_slide)
    _remap_relationship_ids(duplicated_slide._element, relationship_id_map)
    return duplicated_slide


def _remap_relationship_ids(slide_element, relationship_id_map: dict[str, str]) -> None:
    if not relationship_id_map:
        return

    relationship_attributes = (qn("r:embed"), qn("r:link"), qn("r:id"))
    for element in slide_element.iter():
        for attribute_name in relationship_attributes:
            old_relationship_id = element.get(attribute_name)
            if old_relationship_id in relationship_id_map:
                element.set(attribute_name, relationship_id_map[old_relationship_id])


def _set_shape_text(slide, shape_name: str, text: str) -> None:
    for shape in slide.shapes:
        if shape.name == shape_name and hasattr(shape, "text_frame"):
            _replace_text_preserving_template_format(shape, text)
            return
    raise ValueError(f"Could not find template shape '{shape_name}' in the PowerPoint template.")


def _insert_background_picture(slide, presentation: Presentation, image_path: Path) -> None:
    picture = slide.shapes.add_picture(
        str(image_path),
        0,
        0,
        width=presentation.slide_width,
        height=presentation.slide_height,
    )
    picture_element = picture.element
    sp_tree = slide.shapes._spTree
    sp_tree.remove(picture_element)
    sp_tree.insert(2, picture_element)


def _replace_text_preserving_template_format(shape, text: str) -> None:
    text_frame = shape.text_frame
    paragraphs = list(text_frame.paragraphs)
    incoming_lines = (text or "").splitlines() or [""]

    template_paragraph = paragraphs[0]
    _set_paragraph_text(template_paragraph, incoming_lines[0])

    for extra_paragraph in paragraphs[1:]:
        p_element = extra_paragraph._p
        p_element.getparent().remove(p_element)

    for extra_line in incoming_lines[1:]:
        new_paragraph = text_frame.add_paragraph()
        _copy_paragraph_properties(new_paragraph, template_paragraph)
        _set_paragraph_text(new_paragraph, extra_line, template_paragraph)


def _set_paragraph_text(paragraph, text: str, template_paragraph=None) -> None:
    template_paragraph = template_paragraph or paragraph
    runs = list(paragraph.runs)
    template_runs = list(template_paragraph.runs) or runs

    if not runs:
        run = paragraph.add_run()
        if template_runs:
            _copy_run_properties(run, template_runs[0])
        runs = [run]

    runs[0].text = text or " "
    if template_runs and template_runs[0]._r.rPr is not None:
        _copy_run_properties(runs[0], template_runs[0])

    for extra_run in runs[1:]:
        extra_run.text = ""


def _copy_run_properties(run, template_run) -> None:
    template_rpr = template_run._r.rPr
    if template_rpr is None:
        return

    current_rpr = run._r.rPr
    if current_rpr is not None:
        run._r.remove(current_rpr)

    run._r.insert(0, deepcopy(template_rpr))


def _copy_paragraph_properties(paragraph, template_paragraph) -> None:
    template_ppr = template_paragraph._p.pPr
    if template_ppr is None:
        return

    current_ppr = paragraph._p.pPr
    if current_ppr is not None:
        paragraph._p.remove(current_ppr)

    paragraph._p.insert(0, deepcopy(template_ppr))


def _copy_slide_background(source_slide, duplicated_slide) -> None:
    source_bg = source_slide._element.cSld.bg
    if source_bg is None:
        return

    current_bg = duplicated_slide._element.cSld.bg
    if current_bg is not None:
        duplicated_slide._element.cSld.remove(current_bg)

    duplicated_slide._element.cSld.insert(0, deepcopy(source_bg))


def _apply_font_override_to_shape(slide, shape_name: str, font_name: str) -> None:
    for shape in slide.shapes:
        if shape.name != shape_name or not hasattr(shape, "text_frame"):
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text:
                    run.font.name = font_name
        return


def _apply_color_override_to_shape(slide, shape_name: str, color_value: str) -> None:
    rgb = _parse_hex_color(color_value)
    for shape in slide.shapes:
        if shape.name != shape_name or not hasattr(shape, "text_frame"):
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text:
                    run.font.color.rgb = rgb
        return


def _parse_hex_color(color_value: str) -> RGBColor:
    normalized = (color_value or "").strip().lstrip("#")
    if len(normalized) != 6:
        raise ValueError(f"Invalid color value: {color_value}")
    return RGBColor.from_string(normalized.upper())


def _clear_existing_slides(presentation: Presentation) -> None:
    while len(presentation.slides) > 0:
        slide_id = presentation.slides._sldIdLst[0]
        relationship_id = slide_id.rId
        presentation.part.drop_rel(relationship_id)
        del presentation.slides._sldIdLst[0]


def _apply_background(
    slide,
    accent_color: RGBColor,
    template_ppt_path: Path | None = None,
    background_image_path: Path | None = None,
) -> None:
    if background_image_path is not None:
        slide.shapes.add_picture(str(background_image_path), 0, 0, width=SLIDE_WIDTH, height=SLIDE_HEIGHT)

    if template_ppt_path is None:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLOR_BG

    overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = COLOR_BG
    overlay.fill.transparency = 0.18 if background_image_path is not None else (0.28 if template_ppt_path is not None else 0)
    overlay.line.fill.background()

    top_band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.45))
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = COLOR_BG_ALT
    top_band.fill.transparency = 0.12 if (background_image_path is not None or template_ppt_path is not None) else 0
    top_band.line.fill.background()

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(0.55), Inches(2.1), Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()


def _add_cover_slide(
    presentation: Presentation,
    eyebrow: str,
    title: str,
    subtitle: str,
    accent_color: RGBColor,
    template_ppt_path: Path | None = None,
    background_image_path: Path | None = None,
    header_font_name: str | None = None,
    header_font_color: str | None = None,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _apply_background(
        slide,
        accent_color,
        template_ppt_path=template_ppt_path,
        background_image_path=background_image_path,
    )

    eyebrow_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.05), Inches(5.5), Inches(0.5))
    _style_textbox(eyebrow_box, eyebrow, 16, True, COLOR_MUTED, PP_ALIGN.LEFT)

    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.65), Inches(11.2), Inches(1.8))
    _style_textbox(
        title_box,
        title,
        28,
        True,
        _parse_hex_color(header_font_color) if header_font_color else COLOR_TEXT,
        PP_ALIGN.LEFT,
        font_name=header_font_name or "Aptos",
    )

    subtitle_box = slide.shapes.add_textbox(Inches(0.9), Inches(3.2), Inches(10.8), Inches(1.0))
    _style_textbox(
        subtitle_box,
        subtitle,
        18,
        False,
        _parse_hex_color(header_font_color) if header_font_color else COLOR_MUTED,
        PP_ALIGN.LEFT,
        font_name=header_font_name or "Aptos",
    )
    subtitle_box.text_frame.word_wrap = True


def _add_title_slide(presentation: Presentation, title: str, subtitle: str, accent_color: RGBColor) -> None:
    _add_cover_slide(
        presentation,
        eyebrow=subtitle.upper(),
        title=title.upper(),
        subtitle="Prepared with SlideGen",
        accent_color=accent_color,
    )


def _add_lyrics_slide(presentation: Presentation, song_title: str, chunk: SlideChunk) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _apply_background(slide, COLOR_ORANGE if chunk.section_label else COLOR_TEAL)

    meta_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(7.0), Inches(0.35))
    _style_textbox(meta_box, song_title.upper(), 13, True, COLOR_MUTED, PP_ALIGN.LEFT)

    if chunk.section_label:
        label_box = slide.shapes.add_textbox(Inches(10.15), Inches(0.62), Inches(2.1), Inches(0.4))
        _style_textbox(label_box, chunk.section_label, 14, True, COLOR_ORANGE, PP_ALIGN.RIGHT)

    body_box = slide.shapes.add_textbox(Inches(1.1), Inches(1.35), Inches(11.1), Inches(4.95))
    text = "\n".join(chunk.lines)
    _style_textbox(body_box, text, _pick_lyrics_font_size(chunk.lines), True, COLOR_TEXT, PP_ALIGN.CENTER)
    body_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    body_box.text_frame.word_wrap = True


def _add_verse_slide(
    presentation: Presentation,
    reference: str,
    left_label: str,
    right_label: str,
    left_text: str,
    right_text: str,
    template_ppt_path: Path | None = None,
    background_image_path: Path | None = None,
    header_font_name: str | None = None,
    verse_font_name: str | None = None,
    header_font_color: str | None = None,
    verse_font_color: str | None = None,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _apply_background(
        slide,
        COLOR_TEAL,
        template_ppt_path=template_ppt_path,
        background_image_path=background_image_path,
    )

    ref_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.7), Inches(4.4), Inches(0.45))
    _style_textbox(
        ref_box,
        reference.upper(),
        18,
        True,
        _parse_hex_color(header_font_color) if header_font_color else COLOR_TEXT,
        PP_ALIGN.LEFT,
        font_name=header_font_name or "Aptos",
    )

    left_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.45), Inches(5.65), Inches(5.25))
    right_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.45), Inches(5.65), Inches(5.25))

    for panel in (left_panel, right_panel):
        panel.fill.solid()
        panel.fill.fore_color.rgb = COLOR_PANEL
        panel.line.fill.background()

    left_label_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.7), Inches(1.8), Inches(0.35))
    right_label_box = slide.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(1.8), Inches(0.35))
    _style_textbox(
        left_label_box,
        left_label,
        14,
        True,
        _parse_hex_color(header_font_color) if header_font_color else COLOR_TEAL,
        PP_ALIGN.LEFT,
        font_name=header_font_name or "Aptos",
    )
    _style_textbox(
        right_label_box,
        right_label,
        14,
        True,
        _parse_hex_color(header_font_color) if header_font_color else COLOR_ORANGE,
        PP_ALIGN.LEFT,
        font_name=header_font_name or "Aptos",
    )

    left_body = slide.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(5.0), Inches(4.15))
    right_body = slide.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.0), Inches(4.15))
    _style_textbox(
        left_body,
        left_text,
        _pick_verse_font_size(left_text),
        False,
        _parse_hex_color(verse_font_color) if verse_font_color else COLOR_TEXT,
        PP_ALIGN.LEFT,
        font_name=verse_font_name or "Aptos",
    )
    _style_textbox(
        right_body,
        right_text,
        _pick_verse_font_size(right_text),
        False,
        _parse_hex_color(verse_font_color) if verse_font_color else COLOR_TEXT,
        PP_ALIGN.LEFT,
        font_name=verse_font_name or "Aptos",
    )
    left_body.text_frame.word_wrap = True
    right_body.text_frame.word_wrap = True
    left_body.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    right_body.text_frame.vertical_anchor = MSO_ANCHOR.TOP


def _pick_lyrics_font_size(lines: list[str]) -> int:
    longest = max((len(line) for line in lines), default=0)
    line_count = len(lines)
    if line_count <= 2 and longest < 40:
        return 24
    if line_count <= 3 and longest < 56:
        return 21
    return 18


def _pick_verse_font_size(text: str) -> int:
    length = len(text)
    if length < 110:
        return 18
    if length < 180:
        return 16
    return 14


def _style_textbox(shape, text: str, font_size: int, bold: bool, color: RGBColor, alignment, font_name: str = "Aptos") -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = alignment
    run = paragraph.add_run()
    run.text = text or " "
    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color
